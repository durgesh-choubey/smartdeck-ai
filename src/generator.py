"""
SmartDeck AI -- core generation logic.

Two responsibilities, kept deliberately separate:

1. `read_workbook_data()` -- reads the Excel workbook (Info / KPIs /
   Highlights / Chart sheets) into a plain dict. This is the single source
   of truth for "what goes in the deck"; the Streamlit app also uses this
   same dict to drive the in-app preview, so the preview and the generated
   PPTX can never drift out of sync with each other.

2. `generate_pptx()` -- opens the bundled template and fills in the
   placeholders using that dict. Text placeholders are written as
   `run.text = ...` rather than `text_frame.text = ...`, which matters:
   assigning to text_frame.text destroys the run's existing font formatting,
   while assigning to an existing run's .text preserves it. The chart is
   updated via python-pptx's native `chart.replace_data()`, which keeps it a
   real, still-editable PowerPoint chart (not a static image).
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field

import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData


# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------
@dataclass
class DeckData:
    info: dict[str, str] = field(default_factory=dict)
    kpis: list[tuple[str, str]] = field(default_factory=list)
    highlights: list[str] = field(default_factory=list)
    chart_title: str = "Revenue by Region"
    chart_categories: list[str] = field(default_factory=list)
    chart_values: list[float] = field(default_factory=list)

    def footer_info(self) -> str:
        parts = [
            self.info.get("company", ""),
            self.info.get("date", ""),
            self.info.get("prepared_by", ""),
        ]
        return "  |  ".join(p for p in parts if p)

    def kpi_list_text(self) -> str:
        return "\n".join(f"{label}:  {value}" for label, value in self.kpis)

    def highlights_list_text(self) -> str:
        return "\n".join(f"\u2022 {h}" for h in self.highlights)


# ---------------------------------------------------------------------------
# Reading the Excel workbook
# ---------------------------------------------------------------------------
def read_workbook_data(path_or_buffer) -> DeckData:
    wb = openpyxl.load_workbook(path_or_buffer, data_only=True)

    info: dict[str, str] = {}
    if "Info" in wb.sheetnames:
        ws = wb["Info"]
        for row in ws.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                info[str(row[0]).strip()] = "" if row[1] is None else str(row[1]).strip()

    kpis: list[tuple[str, str]] = []
    if "KPIs" in wb.sheetnames:
        ws = wb["KPIs"]
        for row in ws.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                kpis.append((str(row[0]).strip(), "" if row[1] is None else str(row[1]).strip()))

    highlights: list[str] = []
    if "Highlights" in wb.sheetnames:
        ws = wb["Highlights"]
        for row in ws.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                highlights.append(str(row[0]).strip())

    chart_title = "Revenue by Region"
    chart_categories: list[str] = []
    chart_values: list[float] = []
    if "Chart" in wb.sheetnames:
        ws = wb["Chart"]
        first_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), None)
        if first_row and first_row[0] == "chart_title" and first_row[1]:
            chart_title = str(first_row[1]).strip()
        for row in ws.iter_rows(min_row=4, values_only=True):
            if row and row[0] is not None and row[1] is not None:
                try:
                    chart_categories.append(str(row[0]).strip())
                    chart_values.append(float(row[1]))
                except (TypeError, ValueError):
                    continue

    return DeckData(
        info=info, kpis=kpis, highlights=highlights,
        chart_title=chart_title, chart_categories=chart_categories, chart_values=chart_values,
    )


# ---------------------------------------------------------------------------
# Filling the PPTX template
# ---------------------------------------------------------------------------
_TEXT_REPLACEMENTS_SINGLE_LINE = {
    "{{footer_info}}": lambda d: d.footer_info(),
    "{{chart_title}}": lambda d: d.chart_title,
}

_TEXT_REPLACEMENTS_TITLE = {
    "{{report_title}}": lambda d: d.info.get("report_title", ""),
    "{{report_subtitle}}": lambda d: d.info.get("report_subtitle", ""),
}


def _set_run_text_multiline(paragraph_shape_text_frame, text: str) -> None:
    """Replace a single-run placeholder paragraph's text with possibly
    multi-line content, preserving the first run's font formatting for
    every resulting line."""
    tf = paragraph_shape_text_frame
    first_para = tf.paragraphs[0]
    template_run = first_para.runs[0]
    font = template_run.font
    lines = text.split("\n") or [""]

    template_run.text = lines[0]

    # Clear any extra paragraphs left over from the template, then add one
    # new paragraph per remaining line, copying the template run's font.
    for extra_para in list(tf.paragraphs[1:]):
        extra_para._p.getparent().remove(extra_para._p)

    for line in lines[1:]:
        new_para = tf.add_paragraph()
        run = new_para.add_run()
        run.text = line
        run.font.size = font.size
        run.font.bold = font.bold
        run.font.italic = font.italic
        if font.color and font.color.type is not None:
            run.font.color.rgb = font.color.rgb


def generate_pptx(template_path: str, data: DeckData) -> bytes:
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        marker = run.text.strip()
                        if marker in _TEXT_REPLACEMENTS_SINGLE_LINE:
                            run.text = _TEXT_REPLACEMENTS_SINGLE_LINE[marker](data)
                        elif marker in _TEXT_REPLACEMENTS_TITLE:
                            run.text = _TEXT_REPLACEMENTS_TITLE[marker](data)
                        elif marker == "{{kpi_list}}":
                            _set_run_text_multiline(shape.text_frame, data.kpi_list_text())
                        elif marker == "{{highlights_list}}":
                            _set_run_text_multiline(shape.text_frame, data.highlights_list_text())

            if shape.has_chart:
                if data.chart_categories:
                    chart_data = CategoryChartData()
                    chart_data.categories = data.chart_categories
                    chart_data.add_series("Revenue ($K)", tuple(data.chart_values))
                    shape.chart.replace_data(chart_data)
                else:
                    # Blank state: replace the template's sample numbers with an
                    # explicit "no data yet" placeholder rather than leaving the
                    # old sample chart in place.
                    chart_data = CategoryChartData()
                    chart_data.categories = ["No data yet"]
                    chart_data.add_series("Revenue ($K)", (0,))
                    shape.chart.replace_data(chart_data)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()